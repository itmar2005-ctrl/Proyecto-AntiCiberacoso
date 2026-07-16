import os
import uuid
from typing import Optional
from app.config import settings

class DocumentService:
    def __init__(self):
        os.makedirs(settings.upload_dir, exist_ok=True)

    def extract_text(self, filepath: str) -> Optional[str]:
        ext = os.path.splitext(filepath)[1].lower()
        try:
            if ext == ".txt":
                with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                    return f.read()
            elif ext == ".pdf":
                return self._extract_pdf(filepath)
            elif ext in (".docx", ".doc"):
                return self._extract_docx(filepath)
            elif ext == ".pptx":
                return self._extract_pptx(filepath)
            elif ext in (".csv", ".tsv"):
                return self._extract_csv(filepath)
            elif ext in (".xlsx", ".xls"):
                return self._extract_excel(filepath)
            elif ext in (".html", ".htm"):
                return self._extract_html(filepath)
            elif ext == ".json":
                return self._extract_json(filepath)
            else:
                return None
        except Exception as e:
            return f"[Error al extraer texto: {e}]"

    def _extract_pdf(self, path: str) -> str:
        from pypdf import PdfReader
        reader = PdfReader(path)
        return "\n".join([page.extract_text() for page in reader.pages])

    def _extract_docx(self, path: str) -> str:
        try:
            from docx import Document
            doc = Document(path)
            return "\n".join([p.text for p in doc.paragraphs])
        except ImportError:
            return "[DOCX: instala python-docx]"

    def _extract_pptx(self, path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        except ImportError:
            return "[PPTX: instala python-pptx]"

    def _extract_csv(self, path: str) -> str:
        import csv
        rows = []
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(", ".join(row))
        return "\n".join(rows)

    def _extract_excel(self, path: str) -> str:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            texts = []
            for sheet in wb.worksheets:
                texts.append(f"--- Hoja: {sheet.title} ---")
                for row in sheet.iter_rows(values_only=True):
                    row_text = ", ".join([str(c) if c is not None else "" for c in row])
                    if row_text.strip():
                        texts.append(row_text)
            return "\n".join(texts)
        except ImportError:
            return "[Excel: instala openpyxl]"

    def _extract_html(self, path: str) -> str:
        try:
            from bs4 import BeautifulSoup
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            return soup.get_text(separator="\n")
        except ImportError:
            return "[HTML: instala beautifulsoup4]"

    def _extract_json(self, path: str) -> str:
        import json
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, indent=2, ensure_ascii=False)

    def save_file(self, filename: str, content: bytes) -> tuple[str, str]:
        doc_id = str(uuid.uuid4())
        safe_name = f"{doc_id}_{filename}"
        filepath = os.path.join(settings.upload_dir, safe_name)
        with open(filepath, "wb") as f:
            f.write(content)
        return doc_id, filepath
